"""Fix PPTX to add actual hyperlinks."""
from pptx import Presentation
from pptx.util import Inches, Pt

# Load existing PPTX
prs = Presentation("KnowledgeWorker-E2E-Evidence.pptx")

# Add hyperlinks to slides that reference GitHub
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text
                if "github.com" in text.lower() or ".py" in text or ".md" in text:
                    # Extract URL if present
                    if "http" in text:
                        url_start = text.find("http")
                        url_end = text.find(" ", url_start)
                        if url_end == -1:
                            url_end = len(text)
                        url = text[url_start:url_end]
                        
                        # Add hyperlink
                        for run in paragraph.runs:
                            if url in run.text:
                                run.hyperlink.address = url
                                print(f"✓ Added hyperlink on slide {slide_idx + 1}: {url[:60]}...")

# Save
prs.save("KnowledgeWorker-E2E-Evidence.pptx")
print(f"\n✅ PPTX hyperlinks fixed!")
