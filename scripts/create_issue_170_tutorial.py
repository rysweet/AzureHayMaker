#!/usr/bin/env python3
"""Create PowerPoint tutorial for Issue #170 - Graph SDK Bug Fix & Deployment"""

from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    content = slide.placeholders[1].text_frame
    content.clear()

    for bullet in bullets:
        p = content.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Issue #170: Graph SDK Bug Fix",
        "Complete E2E Knowledge Worker Deployment\nDecember 12, 2025",
    )

    # Slide 2: The Bug
    add_content_slide(
        prs,
        "The Blocking Bug",
        [
            "Error: 'dict' object has no attribute 'headers'",
            "Location: permission_granter.py line 72",
            "Impact: Prevented Knowledge Worker deployment",
            "Root Cause: Passing dict instead of RequestConfiguration object",
        ],
    )

    # Slide 3: The Fix
    add_content_slide(
        prs,
        "The Solution",
        [
            "Use proper RequestConfiguration from kiota_abstractions",
            "Match pattern from sp_manager.py:639-648",
            "Add OData injection prevention (_sanitize_odata_value)",
            "Update all tests to work with new pattern",
            "PR #171: https://github.com/rysweet/AzureHayMaker/pull/171",
        ],
    )

    # Slide 4: Code Changes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Code Fix Details"

    tf = slide.placeholders[1].text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "BEFORE (Broken):"
    p.font.bold = True
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = 'request_configuration={"query_parameters": {"filter": "..."}}'
    p.font.name = "Courier New"
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "AFTER (Fixed):"
    p.font.bold = True
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "request_config = RequestConfiguration()"
    p.font.name = "Courier New"
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "request_config.query_parameters = query_params"
    p.font.name = "Courier New"
    p.font.size = Pt(14)

    # Slide 5: Permission Setup
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "API Permissions Required"

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for perm in [
        "Directory.Read.All - Read service principals",
        "User.ReadWrite.All - Create users",
        "AppRoleAssignment.ReadWrite.All - Grant permissions",
        "Mail.ReadWrite - Email operations",
    ]:
        p = tf.add_paragraph()
        p.text = perm
        p.level = 0
        p.font.size = Pt(18)

    # Slide 6: Grant Permissions Command
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "How to Grant Permissions"

    tf = slide.placeholders[1].text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Using Azure CLI with admin rights:"
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""

    cmd_text = """az rest --method POST \\
  --uri "https://graph.microsoft.com/v1.0/\\
    servicePrincipals/{GRAPH_SP}/appRoleAssignedTo" \\
  --body '{
    "principalId": "{YOUR_SP}",
    "resourceId": "{GRAPH_SP}",
    "appRoleId": "7ab1d382-f21e-4acd-a863-ba3e13f7da61"
  }' """

    p = tf.add_paragraph()
    p.text = cmd_text
    p.font.name = "Courier New"
    p.font.size = Pt(12)

    # Slide 7: Deployment Results
    add_content_slide(
        prs,
        "Deployment Results",
        [
            "✅ All 5 workers created successfully",
            "✅ No Graph SDK errors",
            "✅ No permission errors",
            "⚠️ E5 licenses exhausted (25/25 used)",
            "⚠️ Mailbox provisioning takes 15+ minutes",
            "⚠️ Anthropic API 500 errors (service issue)",
        ],
    )

    # Slide 8: Test Results
    add_content_slide(
        prs,
        "Test & CI Results",
        [
            "41 tests passing (permission_granter + knowledge_worker)",
            "All CI checks green ✅",
            "0 linting errors",
            "Fixed async test markers (@pytest.mark.anyio)",
            "Fixed import sorting and whitespace",
        ],
    )

    # Slide 9: Key Learnings
    add_content_slide(
        prs,
        "Key Learnings",
        [
            "Microsoft Graph SDK requires proper RequestConfiguration objects",
            "Inline imports needed ruff import sorting fixes",
            "pytest-anyio needs anyio_backend fixture for asyncio-only",
            "Permission grants via az rest API work reliably",
            "Exchange Online mailbox provisioning is slow (15+ min)",
        ],
    )

    # Slide 10: Success Metrics
    add_content_slide(
        prs,
        "Success Metrics",
        [
            "PR #171: 4 commits, 9 files changed",
            "Security: OData injection prevention added",
            "Coverage: 89% on permission_granter.py",
            "Duration: Issue opened to resolution in < 24 hours",
            "Deployment: 5/5 workers created successfully",
        ],
    )

    # Save
    output_file = "/home/azureuser/src/AzureHayMaker/Issue_170_Tutorial.pptx"
    prs.save(output_file)
    print(f"✅ PowerPoint created: {output_file}")


if __name__ == "__main__":
    main()
