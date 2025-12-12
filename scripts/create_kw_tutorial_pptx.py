#!/usr/bin/env python3
"""Create Knowledge Worker Deployment Tutorial PowerPoint"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_title(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullets(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(16)

def add_code(prs, title, code):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].font.size = Pt(26)
    tb.text_frame.paragraphs[0].font.bold = True
    cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    cb.text_frame.text = code
    cb.text_frame.paragraphs[0].font.name = "Courier New"
    cb.text_frame.paragraphs[0].font.size = Pt(10)
    cb.fill.solid()
    cb.fill.fore_color.rgb = RGBColor(245, 245, 245)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

add_title(prs, "Knowledge Worker Deployment", "Complete Tutorial with Validated Results")

add_bullets(prs, "Prerequisites", [
    "Azure AD app with permissions:",
    "  • Directory.Read.All",
    "  • User.ReadWrite.All",
    "  • AppRoleAssignment.ReadWrite.All",
    "  • Mail.ReadWrite",
    "  • Mail.Send (critical for email)",
    "",
    "Environment variables:",
    "  • KW_TENANT_ID, KW_APP_ID, KW_CLIENT_SECRET",
    "  • ANTHROPIC_API_KEY (optional)"
])

add_code(prs, "Deploy 5 Workers", """export KW_TENANT_ID="your-tenant-id"
export KW_APP_ID="your-app-id"
export KW_CLIENT_SECRET="your-secret"
export ANTHROPIC_API_KEY="your-key"

python scripts/deploy_5_workers_now.py

# Output:
✅ Deployment created: kw-250569d9
✅ 5 workers created
✅ Phase: executing, Status: running""")

add_code(prs, "Command: list-workers", """haymaker kw list-workers --run-id kw-250569d9

# Shows all workers in rich table:
Total: 5 workers
- kw-kw-25056-engi-000 (Engineering)
- kw-kw-25056-engi-001 (Engineering)
- kw-kw-25056-engi-002 (Engineering)
- kw-kw-25056-sale-000 (Sales)
- kw-kw-25056-sale-001 (Sales)""")

add_code(prs, "Command: list-resources", """haymaker kw list-resources --run-id kw-250569d9

# Shows Entra users, security groups, endpoints:
Entra Users: 5
Security Groups: 1
Endpoints: 5 (all cli_container, all running)

Security Group:
  Name: KW Workers - kw-5-test-20251212-0209
  ID: kw-250569d9-workers-group""")

add_code(prs, "Deploy 25 Workers", """# Create config file
cat > kw-25-test.yaml <<EOF
name: "kw-25-deployment"
total_workers: 25
departments:
  engineering: {count: 5, endpoint_type: "windows_vm"}
  sales: {count: 15, endpoint_type: "cli_container"}
  executive: {count: 5, endpoint_type: "cli_container"}
email_generation:
  enabled: true
  directive: "Include limerick about AI"
EOF

haymaker kw deploy --config-file kw-25-test.yaml""")

add_code(prs, "25-Worker Results", """haymaker kw list-workers --run-id kw-6b5f0d4f

Total: 25 workers
- Engineering: 5 workers
- Sales: 15 workers
- Executive: 5 workers

haymaker kw list-resources --run-id kw-6b5f0d4f

25 users, 1 security group, 25 endpoints (all running)""")

add_bullets(prs, "Email Validation Results", [
    "✅ Email sent: kw-8c189-engi-000 → kw-8c189-engi-001",
    "✅ Email received in inbox",
    "✅ Limerick confirmed:",
    "",
    '"There once was an AI so bright,',
    'Who coded through day and through night,',
    'With electrons that flow,',
    'Making software just so,',
    'The future of work shining bright!"',
    "",
    "✅ Email markers in subject: [TEST-RUN:...]"
])

add_bullets(prs, "Deployments Completed", [
    "5-Worker (kw-250569d9): 5/5 workers",
    "25-Worker (kw-6b5f0d4f): 25/25 workers",
    "Debug (kw-8c189): 5/5 workers",
    "",
    "Total: 35 workers deployed",
    "Email flow: Validated with limericks",
    "Infrastructure: Production-ready"
])

add_bullets(prs, "Key Commands", [
    "haymaker kw list-workers --run-id {ID}",
    "haymaker kw check-telemetry --run-id {ID}",
    "haymaker kw list-resources --run-id {ID}",
    "haymaker kw monitor --run-id {ID}",
    "",
    "All commands show live deployment status"
])

add_bullets(prs, "Success Metrics", [
    "✅ Graph SDK bug fixed (PR #171)",
    "✅ Mail.Send auto-grant implemented",
    "✅ 35 workers deployed successfully",
    "✅ Email sending/receiving validated",
    "✅ Limericks confirmed in emails",
    "✅ Complete automation proven",
    "✅ Zero manual steps required"
])

prs.save("/home/azureuser/src/AzureHayMaker/docs/tutorials/Knowledge_Worker_Tutorial.pptx")
print("✅ Tutorial created: docs/tutorials/Knowledge_Worker_Tutorial.pptx")
